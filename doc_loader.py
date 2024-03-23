from dotenv import load_dotenv
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
import os
from PyPDF2 import PdfReader
from langchain.docstore.document import Document
import google.generativeai as genai 
import csv
import requests
from bs4 import BeautifulSoup
from pptx import Presentation


class DocumentChunker:
    '''
    a collection of all different file types that can be chunked.
    Current file types supported: [PDF]
    
    '''
    def __init__(self,files,course_name):
        load_dotenv()
        genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
        #pc = pinecone.Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
        self.files = files # files : file_id, name, type, size
        self.course_name = course_name

    def read_files(self):
        text=""
        docs = [] # List of Documents
        for file in self.files:
            file_type = os.path.splitext(file.name)[1]
            if file_type == ".pdf":
                docs.append(Document(page_content=self.pdf_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}","course_id": f"{self.course_name}"}))
            elif file_type == ".csv":
                docs.append(Document(page_content=self.csv_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}","course_id": f"{self.course_name}"}))
            elif file_type == ".pptx":
                docs.append(Document(page_content=self.ppt_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}","course_id": f"{self.course_name}"}))
            elif file_type.startswith("http"):  # A simple check, you might want to refine this.
                  docs.append(Document(page_content=self.url_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}","course_id": f"{self.course_name}"}))
            elif file_type == ".txt":
                  docs.append(Document(page_content=self.txt_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}","course_id": f"{self.course_name}"}))
            '''
            if file_type == ".csv":
                docs.append(Document(page_content=self.csv_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}"}))
            if file_type == ".txt":
                docs.append(Document(page_content=self.txt_reader(file),metadata={"source": f"{file.name}", "type": f"{file_type}"})) 
            '''
        chunked_docs = self.get_chunks(docs)             
        return chunked_docs
        #vector_store = self.add_vector_to_store(docs)




    def pdf_reader(self,file):
        text=""
        docs = []
        pdf_reader=PdfReader(file)
        for page in pdf_reader.pages:
            text+=page.extract_text()
        #docs.append(Document(page_content=text, metadata={"source": f"{file.name}", "type": f"{file.type}"}))
        return text

    def csv_reader(self, file):
        with open(file, mode='r', encoding='utf-8') as csv_file:
            csv_reader = csv.reader(csv_file)
            text = " ".join([" ".join(row) for row in csv_reader])
        return text

    def ppt_reader(self, file):
        from io import BytesIO
        bytes_data = file.getvalue()
        prs = Presentation(BytesIO(bytes_data))
        text = ""
        for slide in prs.slides:
            if slide.shapes.title:
                text +="Topic:"+slide.shapes.title.text+"-"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    text += shape.text + " "

        return text

    def url_reader(self, url):
        response = requests.get(url.name)
        soup = BeautifulSoup(response.text, 'html.parser')
        text = soup.get_text()
        return text

    def txt_reader(self, file):
        if file.type=='text/plain':
            from io import StringIO
            stringio=StringIO(file.getvalue().decode('utf-8'))
            text=stringio.read()
        return text
    
    def unstructured_reader(self,file):
        pass

    # Using a standerd RecursiveCharacterTextSplitter
    def get_chunks(self,docs):
        text_splitter=RecursiveCharacterTextSplitter(chunk_size=2000,chunk_overlap=200)
        chunks = text_splitter.split_documents(docs)            
        return chunks
    
